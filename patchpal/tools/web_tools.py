"""Web tools (fetch, search)."""

import os

import requests
from bs4 import BeautifulSoup

try:
    from ddgs import DDGS
except ImportError:
    from duckduckgo_search import DDGS

try:
    import pymupdf

    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    import docx

    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    import pptx

    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

from patchpal.tools.common import (
    MAX_WEB_CONTENT_CHARS,
    MAX_WEB_CONTENT_SIZE,
    WEB_HEADERS,
    WEB_REQUEST_TIMEOUT,
    _get_permission_manager,
    _operation_limiter,
    audit_logger,
)


def web_fetch(url: str, extract_text: bool = True) -> str:
    """
    Fetch content from a URL and optionally extract readable text.

    Args:
        url: The URL to fetch
        extract_text: If True, extract readable text from HTML/PDF (default: True)

    Returns:
        The fetched content (text extracted from HTML/PDF if extract_text=True)

    Raises:
        ValueError: If request fails or content is too large
    """
    # Check permission before proceeding
    permission_manager = _get_permission_manager()
    description = f"   Fetch: {url}"
    if not permission_manager.request_permission("web_fetch", description):
        return "Operation cancelled by user."

    _operation_limiter.check_limit(f"web_fetch({url[:50]}...)")

    # Validate URL format
    if not url.startswith(("http://", "https://")):
        raise ValueError("URL must start with http:// or https://")

    try:
        # Make request with timeout and browser-like headers
        response = requests.get(
            url,
            timeout=WEB_REQUEST_TIMEOUT,
            headers=WEB_HEADERS,
            stream=True,  # Stream to check size first
            allow_redirects=True,  # Follow redirects (including moved repos)
        )
        response.raise_for_status()

        # Check content size
        content_length = response.headers.get("Content-Length")
        if content_length and int(content_length) > MAX_WEB_CONTENT_SIZE:
            raise ValueError(
                f"Content too large: {int(content_length):,} bytes "
                f"(max {MAX_WEB_CONTENT_SIZE:,} bytes)"
            )

        # Read content with size limit
        content = b""
        for chunk in response.iter_content(chunk_size=8192):
            content += chunk
            if len(content) > MAX_WEB_CONTENT_SIZE:
                raise ValueError(f"Content exceeds size limit ({MAX_WEB_CONTENT_SIZE:,} bytes)")

        # Get content type
        content_type = response.headers.get("Content-Type", "").lower()

        # Extract text based on content type
        if extract_text:
            if "pdf" in content_type and PYMUPDF_AVAILABLE:
                # Extract text from PDF
                try:
                    pdf_document = pymupdf.open(stream=content, filetype="pdf")
                    text_parts = []
                    for page_num in range(pdf_document.page_count):
                        page = pdf_document[page_num]
                        text_parts.append(page.get_text())
                    pdf_document.close()
                    text_content = "\n".join(text_parts)
                except Exception as pdf_error:
                    # Fall back to showing error if PDF extraction fails
                    text_content = (
                        f"[PDF extraction failed: {pdf_error}]\n\n"
                        f"PDF URL: {url}\n"
                        f"To read this PDF, download it locally or try a different source."
                    )
            elif (
                "wordprocessingml" in content_type or "msword" in content_type
            ) and PYTHON_DOCX_AVAILABLE:
                # Extract text from DOCX (or DOC if saved as docx)
                try:
                    import io

                    doc = docx.Document(io.BytesIO(content))
                    text_parts = []
                    for paragraph in doc.paragraphs:
                        text_parts.append(paragraph.text)
                    text_content = "\n".join(text_parts)
                except Exception as docx_error:
                    text_content = (
                        f"[DOCX extraction failed: {docx_error}]\n\n"
                        f"Content-Type: {content_type}\n"
                        f"URL: {url}\n"
                        f"To read this document, download it locally or try a different source."
                    )
            elif (
                "presentationml" in content_type or "ms-powerpoint" in content_type
            ) and PYTHON_PPTX_AVAILABLE:
                # Extract text from PPTX (or PPT if saved as pptx)
                try:
                    import io

                    prs = pptx.Presentation(io.BytesIO(content))
                    text_parts = []
                    for slide_num, slide in enumerate(prs.slides, 1):
                        text_parts.append(f"\n--- Slide {slide_num} ---")
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_parts.append(shape.text)
                    text_content = "\n".join(text_parts)
                except Exception as pptx_error:
                    text_content = (
                        f"[PPTX extraction failed: {pptx_error}]\n\n"
                        f"Content-Type: {content_type}\n"
                        f"URL: {url}\n"
                        f"To read this presentation, download it locally or try a different source."
                    )
            elif "html" in content_type:
                # Extract text from HTML
                text_content = content.decode(response.encoding or "utf-8", errors="replace")
                soup = BeautifulSoup(text_content, "html.parser")

                # Remove script and style elements
                for element in soup(["script", "style", "nav", "footer", "header"]):
                    element.decompose()

                # Get text
                text = soup.get_text()

                # Clean up whitespace
                lines = (line.strip() for line in text.splitlines())
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                text_content = "\n".join(chunk for chunk in chunks if chunk)
            else:
                # For other content types, check if it's a known binary format
                binary_formats = [
                    "image/",
                    "video/",
                    "audio/",
                    "application/zip",
                    "application/x-zip",
                    "application/x-rar",
                    "application/x-tar",
                    "spreadsheetml",  # Excel files (xlsx) - not yet supported
                    "ms-excel",  # Legacy Excel files (xls) - not yet supported
                    "application/octet-stream",
                ]
                is_binary = any(fmt in content_type for fmt in binary_formats)

                if is_binary:
                    text_content = (
                        f"[WARNING: Unsupported binary format]\n\n"
                        f"Content-Type: {content_type}\n"
                        f"URL: {url}\n\n"
                        f"This appears to be a binary file format that cannot be extracted as text.\n"
                        f"Supported formats: HTML, PDF, DOCX, PPTX, plain text, JSON, XML.\n"
                        f"To access this content, download it locally or use a format-specific tool."
                    )
                else:
                    # Assume it's text-based (JSON, XML, CSV, etc.)
                    text_content = content.decode(response.encoding or "utf-8", errors="replace")
        else:
            # No text extraction - just decode
            text_content = content.decode(response.encoding or "utf-8", errors="replace")

        # Truncate content if it exceeds character limit to prevent context window overflow
        if len(text_content) > MAX_WEB_CONTENT_CHARS:
            truncated_content = text_content[:MAX_WEB_CONTENT_CHARS]
            warning_msg = (
                f"\n\n[WARNING: Content truncated from {len(text_content):,} to "
                f"{MAX_WEB_CONTENT_CHARS:,} characters to prevent context window overflow. "
                f"Set PATCHPAL_MAX_WEB_CHARS environment variable to adjust limit.]"
            )
            audit_logger.info(
                f"WEB_FETCH: {url} ({len(text_content)} chars, truncated to {MAX_WEB_CONTENT_CHARS})"
            )
            return truncated_content + warning_msg

        audit_logger.info(f"WEB_FETCH: {url} ({len(text_content)} chars)")
        return text_content

    except requests.Timeout:
        raise ValueError(f"Request timed out after {WEB_REQUEST_TIMEOUT} seconds")
    except requests.RequestException as e:
        raise ValueError(f"Failed to fetch URL: {e}")
    except Exception as e:
        raise ValueError(f"Error processing content: {e}")


def web_search(query: str, max_results: int = 5) -> str:
    """
    Search the web using DuckDuckGo and return results.

    Args:
        query: The search query
        max_results: Maximum number of results to return (default: 5, max: 10)

    Returns:
        Formatted search results with titles, URLs, and snippets

    Raises:
        ValueError: If search fails
    """
    # Check permission before proceeding
    permission_manager = _get_permission_manager()
    description = f"   Search: {query}"
    if not permission_manager.request_permission("web_search", description):
        return "Operation cancelled by user."

    _operation_limiter.check_limit(f"web_search({query[:30]}...)")

    # Limit max_results
    max_results = min(max_results, 10)

    try:
        # Determine SSL verification setting
        # Priority: PATCHPAL_VERIFY_SSL env var > SSL_CERT_FILE > REQUESTS_CA_BUNDLE > default True
        verify_ssl = os.getenv("PATCHPAL_VERIFY_SSL")
        if verify_ssl is not None:
            # User explicitly set PATCHPAL_VERIFY_SSL
            if verify_ssl.lower() in ("false", "0", "no"):
                verify = False
            elif verify_ssl.lower() in ("true", "1", "yes"):
                verify = True
            else:
                # Treat as path to CA bundle
                verify = verify_ssl
        else:
            # Use SSL_CERT_FILE or REQUESTS_CA_BUNDLE if set (for corporate environments)
            verify = os.getenv("SSL_CERT_FILE") or os.getenv("REQUESTS_CA_BUNDLE") or True

        # Perform search using DuckDuckGo
        with DDGS(verify=verify) as ddgs:
            results = list(ddgs.text(query, max_results=max_results))

        if not results:
            audit_logger.info(f"WEB_SEARCH: {query} - No results")
            return f"No search results found for: {query}"

        # Format results
        formatted_results = [f"Search results for: {query}\n"]
        for i, result in enumerate(results, 1):
            title = result.get("title", "No title")
            url = result.get("href", "No URL")
            snippet = result.get("body", "No description")

            formatted_results.append(f"\n{i}. {title}\n   URL: {url}\n   {snippet}")

        output = "\n".join(formatted_results)
        audit_logger.info(f"WEB_SEARCH: {query} - Found {len(results)} results")
        return output

    except Exception as e:
        error_msg = str(e)

        # Provide helpful error messages for common issues
        if "CERTIFICATE_VERIFY_FAILED" in error_msg or "TLS handshake failed" in error_msg:
            return (
                "Web search unavailable: SSL certificate verification failed.\n"
                "This may be due to:\n"
                "- Corporate proxy/firewall blocking requests\n"
                "- Network configuration issues\n"
                "- VPN interference\n\n"
                "Consider using web_fetch with a specific URL if you have one."
            )
        elif "RuntimeError" in error_msg or "error sending request" in error_msg:
            return (
                "Web search unavailable: Network connection failed.\n"
                "Please check your internet connection and try again."
            )
        else:
            raise ValueError(f"Web search failed: {e}")
